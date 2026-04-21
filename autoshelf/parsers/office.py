from __future__ import annotations

from pathlib import Path

from autoshelf.parsers.base import ParsedContext, ParserSpec


def parse_office(path: Path, max_head_chars: int = 2000) -> ParsedContext:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return _parse_docx(path, max_head_chars)
    if suffix == ".pptx":
        return _parse_pptx(path, max_head_chars)
    if suffix == ".xlsx":
        return _parse_xlsx(path, max_head_chars)
    if suffix == ".xls":
        return _parse_xls(path, max_head_chars)
    return ParsedContext(title=path.stem, head_text="", extra_meta={})


def _parse_docx(path: Path, max_head_chars: int) -> ParsedContext:
    try:
        import docx  # type: ignore[import-not-found]
    except ImportError:
        return ParsedContext(title=path.stem, head_text="", extra_meta={"parser": "unavailable"})
    try:
        document = docx.Document(str(path))
        paragraphs = [paragraph for paragraph in document.paragraphs if paragraph.text.strip()]
        text = "\n".join(paragraph.text for paragraph in paragraphs[:2])
        head = text[:max_head_chars]
        title = next(
            (
                paragraph.text.strip()
                for paragraph in paragraphs
                if paragraph.style is not None and "title" in paragraph.style.name.lower()
            ),
            paragraphs[0].text.strip() if paragraphs else path.stem,
        )
        headings = [
            paragraph.text.strip()
            for paragraph in paragraphs
            if paragraph.style is not None and "heading" in paragraph.style.name.lower()
        ]
        return ParsedContext(title=title[:120], head_text=head, extra_meta={"headings": headings[:5]})
    except Exception:
        return ParsedContext(title=path.stem, head_text="", extra_meta={"parser": "failed"})


def _parse_pptx(path: Path, max_head_chars: int) -> ParsedContext:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError:
        return ParsedContext(title=path.stem, head_text="", extra_meta={"parser": "unavailable"})
    try:
        presentation = Presentation(str(path))
        texts: list[str] = []
        titles: list[str] = []
        for slide in presentation.slides[:3]:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
                    titles.append(shape.text.splitlines()[0].strip())
                    break
        content = "\n".join(texts)
        title = titles[0] if titles else path.stem
        return ParsedContext(title=title[:120], head_text=content[:max_head_chars], extra_meta={"slides": len(texts)})
    except Exception:
        return ParsedContext(title=path.stem, head_text="", extra_meta={"parser": "failed"})


def _parse_xlsx(path: Path, max_head_chars: int) -> ParsedContext:
    try:
        import openpyxl  # type: ignore[import-not-found]
    except ImportError:
        return ParsedContext(title=path.stem, head_text="", extra_meta={"parser": "unavailable"})
    try:
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
        texts: list[str] = []
        for sheet in workbook.worksheets[:1]:
            texts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(max_row=10, values_only=True):
                texts.append(" | ".join("" if value is None else str(value) for value in row))
        head_text = "\n".join(texts)[:max_head_chars]
        title = workbook.sheetnames[0] if workbook.sheetnames else path.stem
        return ParsedContext(title=title[:120], head_text=head_text, extra_meta={"sheets": workbook.sheetnames})
    except Exception:
        return ParsedContext(title=path.stem, head_text="", extra_meta={"parser": "failed"})


def _parse_xls(path: Path, max_head_chars: int) -> ParsedContext:
    try:
        import xlrd  # type: ignore[import-not-found]
    except ImportError:
        return ParsedContext(title=path.stem, head_text="", extra_meta={"parser": "unavailable"})
    try:
        workbook = xlrd.open_workbook(path)
        texts: list[str] = []
        for sheet in workbook.sheets()[:2]:
            for row_index in range(min(sheet.nrows, 10)):
                texts.append(" | ".join(str(value) for value in sheet.row_values(row_index)))
        head_text = "\n".join(texts)[:max_head_chars]
        return ParsedContext(title=path.stem, head_text=head_text, extra_meta={})
    except Exception:
        return ParsedContext(title=path.stem, head_text="", extra_meta={"parser": "failed"})


PARSER_SPECS = [
    ParserSpec(name="office-docx", suffixes=(".docx",), parse=_parse_docx),
    ParserSpec(name="office-pptx", suffixes=(".pptx",), parse=_parse_pptx),
    ParserSpec(name="office-xlsx", suffixes=(".xlsx",), parse=_parse_xlsx),
    ParserSpec(name="office-xls", suffixes=(".xls",), parse=_parse_xls),
]
